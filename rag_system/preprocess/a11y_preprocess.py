from __future__ import annotations

import argparse
import base64
import dataclasses
import hashlib
import json
import mimetypes
import os
import platform
import re
import shutil
import subprocess
from dataclasses import dataclass
from pathlib import Path
from typing import Any, Dict, List, Literal, Optional, Tuple

import fitz
from openai import OpenAI
from PIL import Image, ImageDraw
from pydantic import BaseModel, Field

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except Exception:  # pragma: no cover
    Presentation = None
    MSO_SHAPE_TYPE = None

DEFAULT_MODEL = os.getenv('VLM_MODEL', 'gpt-4.1-mini')
DEFAULT_TEMPERATURE = float(os.getenv('VLM_TEMPERATURE', '0.0'))
DEFAULT_DPI = int(os.getenv('RENDER_DPI', '200'))
Label = Literal['error_region', 'error_code']


class ImageClassification(BaseModel):
    label: Label


class CodeExtraction(BaseModel):
    code: str = ''


class SlideFields(BaseModel):
    inspection_item: Optional[str] = None
    error_type: Optional[str] = None
    rationale_bullets: List[str] = Field(default_factory=list)
    fix_code: Optional[str] = None
    error_code_text: Optional[str] = None


@dataclass
class ImageAsset:
    page_index: int
    image_index: int
    path: Path
    sha1: str
    ext: str
    bbox_px: Tuple[int, int, int, int]
    label: Optional[Label] = None


@dataclass
class PageRecord:
    page_index: int
    page_image_path: Path
    page_masked_path: Path
    error_region_images: List[Path]
    error_code_images: List[Path]
    inspection_item: Optional[str]
    error_type: Optional[str]
    rationale_bullets: List[str]
    fix_code: Optional[str]
    error_code_text: Optional[str]
    error_code_text_from_page: Optional[str]
    error_code_texts_from_images: List[str]


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def sha1_bytes(data: bytes) -> str:
    return hashlib.sha1(data).hexdigest()


def safe_json_loads(text: str) -> Any:
    text = text.strip()
    text = re.sub(r'^```(?:json)?\s*', '', text)
    text = re.sub(r'\s*```$', '', text)
    try:
        return json.loads(text)
    except Exception:
        pass
    i, j = text.find('{'), text.rfind('}')
    if i != -1 and j > i:
        return json.loads(text[i:j + 1])
    i, j = text.find('['), text.rfind(']')
    if i != -1 and j > i:
        return json.loads(text[i:j + 1])
    raise ValueError('JSON parse failed')


def _guess_mime(path: Path) -> str:
    mime, _ = mimetypes.guess_type(str(path))
    if mime:
        return mime
    if path.suffix.lower() in {'.jpg', '.jpeg'}:
        return 'image/jpeg'
    return 'image/png'


def _data_url_from_image(path: Path) -> str:
    mime = _guess_mime(path)
    b64 = base64.b64encode(path.read_bytes()).decode('utf-8')
    return f'data:{mime};base64,{b64}'


def _run_cmd(cmd: List[str]) -> None:
    proc = subprocess.run(cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True)
    if proc.returncode != 0:
        raise RuntimeError(f'Command failed: {cmd}\nstdout={proc.stdout}\nstderr={proc.stderr}')


def _find_soffice() -> str:
    for key in ('SOFFICE_PATH', 'LIBREOFFICE_PATH'):
        value = os.getenv(key, '').strip().strip('"')
        if value and Path(value).exists():
            return value
    for candidate in ('soffice', 'soffice.exe'):
        found = shutil.which(candidate)
        if found:
            return found
    system = platform.system().lower()
    candidates: List[str] = []
    if system.startswith('windows'):
        candidates = [
            r'C:\Program Files\LibreOffice\program\soffice.exe',
            r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
        ]
    elif system == 'darwin':
        candidates = ['/Applications/LibreOffice.app/Contents/MacOS/soffice']
    else:
        candidates = ['/usr/bin/soffice', '/usr/local/bin/soffice', '/snap/bin/libreoffice']
    for path in candidates:
        if Path(path).exists():
            return path
    raise RuntimeError('LibreOffice(soffice)를 찾지 못했습니다. PPTX 처리를 위해 LibreOffice 설치 또는 SOFFICE_PATH 지정이 필요합니다.')


class OpenAIVLMBackend:
    def __init__(self, model: str = DEFAULT_MODEL, temperature: float = DEFAULT_TEMPERATURE, api_key: Optional[str] = None):
        api_key = api_key or os.getenv('OPENAI_API_KEY')
        if not api_key:
            raise RuntimeError('OPENAI_API_KEY is required.')
        self.client = OpenAI(api_key=api_key)
        self.model = model
        self.temperature = temperature

    def _responses_parse_available(self) -> bool:
        return hasattr(self.client, 'responses') and hasattr(self.client.responses, 'parse')

    def _schema_json(self, schema_model: Any) -> Dict[str, Any]:
        if hasattr(schema_model, 'model_json_schema'):
            return schema_model.model_json_schema()
        return schema_model.schema()

    def _call_parse_or_fallback(self, *, schema_model: Any, input_payload: List[Dict[str, Any]]) -> Any:
        if self._responses_parse_available():
            try:
                resp = self.client.responses.parse(
                    model=self.model,
                    input=input_payload,
                    temperature=self.temperature,
                    text_format=schema_model,
                )
                return resp.output_parsed
            except Exception:
                pass

        try:
            resp = self.client.responses.create(
                model=self.model,
                input=input_payload,
                temperature=self.temperature,
                text={
                    'format': {
                        'type': 'json_schema',
                        'name': getattr(schema_model, '__name__', 'Schema'),
                        'schema': self._schema_json(schema_model),
                        'strict': True,
                    }
                },
            )
        except TypeError:
            resp = self.client.responses.create(model=self.model, input=input_payload, temperature=self.temperature)
        text = getattr(resp, 'output_text', None) or str(resp)
        return safe_json_loads(text)

    def classify_image(self, image_path: Path) -> Label:
        prompt = (
            '접근성 진단 보고서에서 추출한 이미지를 분류하라. '\
            '반드시 JSON만 출력: {"label":"error_region"} 또는 {"label":"error_code"}. '\
            'error_region은 실제 오류 UI 캡처, error_code는 HTML/CSS/JS/DOM 코드 캡처다.'
        )
        payload = [{
            'role': 'user',
            'content': [
                {'type': 'input_text', 'text': prompt},
                {'type': 'input_image', 'image_url': _data_url_from_image(image_path)},
            ],
        }]
        parsed = self._call_parse_or_fallback(schema_model=ImageClassification, input_payload=payload)
        if isinstance(parsed, ImageClassification):
            return parsed.label
        if isinstance(parsed, dict) and parsed.get('label') in {'error_region', 'error_code'}:
            return parsed['label']
        return 'error_region'

    def extract_code_from_image(self, image_path: Path) -> str:
        prompt = '이미지 안의 코드만 가능한 한 원문 그대로 추출하라. 설명 없이 JSON만 출력: {"code":"..."}'
        payload = [{
            'role': 'user',
            'content': [
                {'type': 'input_text', 'text': prompt},
                {'type': 'input_image', 'image_url': _data_url_from_image(image_path)},
            ],
        }]
        parsed = self._call_parse_or_fallback(schema_model=CodeExtraction, input_payload=payload)
        if isinstance(parsed, CodeExtraction):
            return (parsed.code or '').strip()
        if isinstance(parsed, dict):
            return str(parsed.get('code', '')).strip()
        return ''

    def extract_fields_from_masked_page(self, masked_page_path: Path) -> SlideFields:
        prompt = (
            '접근성 진단 보고서 상세 페이지에서 보이는 텍스트만 추출하라. '
            '마스킹 내부는 무시한다. 반드시 JSON만 출력한다. '
            '{"inspection_item":null,"error_type":null,"rationale_bullets":[],"fix_code":null,"error_code_text":null}'
        )
        payload = [{
            'role': 'user',
            'content': [
                {'type': 'input_text', 'text': prompt},
                {'type': 'input_image', 'image_url': _data_url_from_image(masked_page_path)},
            ],
        }]
        parsed = self._call_parse_or_fallback(schema_model=SlideFields, input_payload=payload)
        if isinstance(parsed, SlideFields):
            return parsed
        if isinstance(parsed, dict):
            return SlideFields(
                inspection_item=parsed.get('inspection_item'),
                error_type=parsed.get('error_type'),
                rationale_bullets=[str(x).strip() for x in (parsed.get('rationale_bullets') or []) if str(x).strip()],
                fix_code=parsed.get('fix_code'),
                error_code_text=parsed.get('error_code_text'),
            )
        return SlideFields()


def _iter_picture_shapes(shapes):
    for shape in shapes:
        try:
            if MSO_SHAPE_TYPE and shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                yield from _iter_picture_shapes(shape.shapes)
            elif MSO_SHAPE_TYPE and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                yield shape
        except Exception:
            continue


def extract_all_images_from_pptx(pptx_path: Path, out_dir: Path):
    if Presentation is None:
        raise RuntimeError('python-pptx가 설치되어 있지 않습니다.')
    prs = Presentation(str(pptx_path))
    slide_w_emu = int(prs.slide_width)
    slide_h_emu = int(prs.slide_height)
    img_root = out_dir / 'extracted_images' / 'pptx'
    ensure_dir(img_root)

    assets: List[ImageAsset] = []
    temp_emu: List[Tuple[int, int, int, int, int, int]] = []
    for si, slide in enumerate(prs.slides):
        slide_dir = img_root / f'slide_{si:03d}'
        ensure_dir(slide_dir)
        pi = 0
        for shape in _iter_picture_shapes(slide.shapes):
            try:
                img = shape.image
                blob = img.blob
                ext = img.ext or 'png'
                p = slide_dir / f'img_{pi:03d}_{sha1_bytes(blob)[:16]}.{ext}'
                p.write_bytes(blob)
                assets.append(ImageAsset(si, pi, p, sha1_bytes(blob), ext, (0, 0, 1, 1), None))
                temp_emu.append((si, pi, int(shape.left), int(shape.top), int(shape.width), int(shape.height)))
                pi += 1
            except Exception:
                continue
    return assets, (slide_w_emu, slide_h_emu), temp_emu


def convert_pptx_to_pdf(pptx_path: Path, out_dir: Path) -> Path:
    ensure_dir(out_dir)
    soffice = _find_soffice()
    _run_cmd([soffice, '--headless', '--convert-to', 'pdf', '--outdir', str(out_dir), str(pptx_path)])
    pdfs = sorted(out_dir.glob('*.pdf'), key=lambda p: p.stat().st_mtime, reverse=True)
    if not pdfs:
        raise FileNotFoundError('PPTX to PDF conversion output not found')
    return pdfs[0]


def render_pdf_to_pngs(pdf_path: Path, out_dir: Path, dpi: int = DEFAULT_DPI) -> List[Path]:
    ensure_dir(out_dir)
    doc = fitz.open(str(pdf_path))
    zoom = dpi / 72.0
    mat = fitz.Matrix(zoom, zoom)
    out: List[Path] = []
    for i in range(doc.page_count):
        pix = doc.load_page(i).get_pixmap(matrix=mat, alpha=False)
        path = out_dir / f'page_{i:03d}.png'
        pix.save(str(path))
        out.append(path)
    doc.close()
    return out


def render_pages_from_pptx(pptx_path: Path, out_dir: Path, dpi: int = DEFAULT_DPI) -> List[Path]:
    render_dir = out_dir / 'rendered_pages' / 'pptx'
    pdf_path = convert_pptx_to_pdf(pptx_path, render_dir)
    return render_pdf_to_pngs(pdf_path, render_dir / 'png', dpi=dpi)


def render_pages_from_pdf(pdf_path: Path, out_dir: Path, dpi: int = DEFAULT_DPI) -> List[Path]:
    return render_pdf_to_pngs(pdf_path, out_dir / 'rendered_pages' / 'pdf' / 'png', dpi=dpi)


def compute_bbox_px_for_pptx_assets(assets: List[ImageAsset], pages_png: List[Path], slide_size_emu: Tuple[int, int], temp_emu: List[Tuple[int, int, int, int, int, int]]) -> None:
    slide_w_emu, slide_h_emu = slide_size_emu
    page_px_size: Dict[int, Tuple[int, int]] = {}
    for idx, path in enumerate(pages_png):
        with Image.open(path) as im:
            page_px_size[idx] = (im.width, im.height)
    asset_map = {(a.page_index, a.image_index): a for a in assets}
    for si, pi, left, top, w_emu, h_emu in temp_emu:
        if si not in page_px_size:
            continue
        pw, ph = page_px_size[si]
        x = int(round(left / slide_w_emu * pw))
        y = int(round(top / slide_h_emu * ph))
        w = int(round(w_emu / slide_w_emu * pw))
        h = int(round(h_emu / slide_h_emu * ph))
        x = max(0, min(x, pw - 1))
        y = max(0, min(y, ph - 1))
        w = max(1, min(w, pw - x))
        h = max(1, min(h, ph - y))
        asset_map[(si, pi)].bbox_px = (x, y, w, h)


def extract_all_images_from_pdf(pdf_path: Path, out_dir: Path, dpi: int = DEFAULT_DPI) -> List[ImageAsset]:
    img_root = out_dir / 'extracted_images' / 'pdf'
    ensure_dir(img_root)
    zoom = dpi / 72.0
    doc = fitz.open(str(pdf_path))
    assets: List[ImageAsset] = []
    for page_i in range(doc.page_count):
        page = doc.load_page(page_i)
        page_dir = img_root / f'page_{page_i:03d}'
        ensure_dir(page_dir)
        img_idx = 0
        for info in page.get_images(full=True):
            xref = info[0]
            rects = page.get_image_rects(xref)
            if not rects:
                continue
            extracted = doc.extract_image(xref)
            img_bytes = extracted.get('image', b'')
            ext = extracted.get('ext', 'png')
            if not img_bytes:
                continue
            img_path = page_dir / f'img_{img_idx:03d}_{sha1_bytes(img_bytes)[:16]}.{ext}'
            if not img_path.exists():
                img_path.write_bytes(img_bytes)
            for rect in rects:
                x0 = int(round(rect.x0 * zoom))
                y0 = int(round(rect.y0 * zoom))
                x1 = int(round(rect.x1 * zoom))
                y1 = int(round(rect.y1 * zoom))
                assets.append(ImageAsset(page_i, img_idx, img_path, sha1_bytes(img_bytes), ext, (x0, y0, max(1, x1 - x0), max(1, y1 - y0)), None))
            img_idx += 1
    doc.close()
    return assets


def mask_page_images(page_png_path: Path, assets_in_page: List[ImageAsset], out_path: Path, pad: int = 2) -> None:
    with Image.open(page_png_path) as im:
        im = im.convert('RGB')
        draw = ImageDraw.Draw(im)
        for asset in assets_in_page:
            x, y, w, h = asset.bbox_px
            draw.rectangle([max(0, x - pad), max(0, y - pad), min(im.width, x + w + pad), min(im.height, y + h + pad)], fill=(255, 255, 255))
        ensure_dir(out_path.parent)
        im.save(out_path)


def run_pipeline(pages_png: List[Path], assets: List[ImageAsset], out_dir: Path, vlm: OpenAIVLMBackend) -> Path:
    for asset in assets:
        asset.label = vlm.classify_image(asset.path)

    assets_by_page: Dict[int, List[ImageAsset]] = {}
    for asset in assets:
        assets_by_page.setdefault(asset.page_index, []).append(asset)

    records: List[PageRecord] = []
    masked_dir = out_dir / 'masked_pages'
    ensure_dir(masked_dir)
    for page_i, page_png in enumerate(pages_png):
        page_assets = assets_by_page.get(page_i, [])
        error_region_imgs = [a.path for a in page_assets if a.label == 'error_region']
        error_code_imgs = [a.path for a in page_assets if a.label == 'error_code']
        masked_path = masked_dir / f'page_{page_i:03d}_masked.png'
        mask_page_images(page_png, page_assets, masked_path)
        fields = vlm.extract_fields_from_masked_page(masked_path)
        code_texts = [txt for txt in (vlm.extract_code_from_image(path) for path in error_code_imgs) if txt]
        error_code_text_from_page = fields.error_code_text.strip() if fields.error_code_text else None
        final_error_code_text = error_code_text_from_page or ('\n\n/* --- error_code_image_split --- */\n\n'.join(code_texts) if code_texts else None)
        records.append(
            PageRecord(
                page_index=page_i,
                page_image_path=page_png,
                page_masked_path=masked_path,
                error_region_images=error_region_imgs,
                error_code_images=error_code_imgs,
                inspection_item=fields.inspection_item.strip() if fields.inspection_item else None,
                error_type=fields.error_type.strip() if fields.error_type else None,
                rationale_bullets=[b.strip() for b in fields.rationale_bullets if b and b.strip()],
                fix_code=fields.fix_code.strip() if fields.fix_code else None,
                error_code_text=final_error_code_text,
                error_code_text_from_page=error_code_text_from_page,
                error_code_texts_from_images=code_texts,
            )
        )

    out_jsonl = out_dir / 'records.jsonl'
    with out_jsonl.open('w', encoding='utf-8') as f:
        for record in records:
            obj = dataclasses.asdict(record)
            obj.pop('page_image_path', None)
            obj.pop('page_masked_path', None)
            obj['error_region_images'] = [str(p) for p in record.error_region_images]
            obj['error_code_images'] = [str(p) for p in record.error_code_images]
            f.write(json.dumps(obj, ensure_ascii=False) + '\n')
    return out_jsonl


def preprocess_input(input_path: Path, out_dir: Path, vlm: OpenAIVLMBackend, dpi: int = DEFAULT_DPI) -> Path:
    ensure_dir(out_dir)
    suffix = input_path.suffix.lower()
    if suffix == '.pptx':
        assets, slide_size_emu, temp_emu = extract_all_images_from_pptx(input_path, out_dir)
        pages_png = render_pages_from_pptx(input_path, out_dir, dpi=dpi)
        compute_bbox_px_for_pptx_assets(assets, pages_png, slide_size_emu, temp_emu)
        return run_pipeline(pages_png, assets, out_dir, vlm)
    if suffix == '.pdf':
        pages_png = render_pages_from_pdf(input_path, out_dir, dpi=dpi)
        assets = extract_all_images_from_pdf(input_path, out_dir, dpi=dpi)
        return run_pipeline(pages_png, assets, out_dir, vlm)
    raise ValueError(f'지원하지 않는 확장자: {suffix}')


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument('--input', required=True)
    parser.add_argument('--out', required=True)
    parser.add_argument('--dpi', type=int, default=DEFAULT_DPI)
    parser.add_argument('--model', default=DEFAULT_MODEL)
    parser.add_argument('--temperature', type=float, default=DEFAULT_TEMPERATURE)
    parser.add_argument('--api-key', default='')
    args = parser.parse_args()
    vlm = OpenAIVLMBackend(model=args.model, temperature=args.temperature, api_key=args.api_key or os.getenv('OPENAI_API_KEY'))
    out_jsonl = preprocess_input(Path(args.input).resolve(), Path(args.out).resolve(), vlm, dpi=args.dpi)
    print(out_jsonl)


if __name__ == '__main__':
    main()
